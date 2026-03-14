"""Parse uploaded files into text chunks with source info (page/slide)."""
from pathlib import Path
import re


def chunk_text(text: str, max_chars: int = 1200, overlap: int = 100) -> list[tuple[str, int | None]]:
    """Split text into overlapping chunks. Returns list of (content, page_or_slide)."""
    if not text.strip():
        return []
    # Prefer splitting on paragraph boundaries
    paragraphs = re.split(r"\n\s*\n", text)
    chunks: list[tuple[str, int | None]] = []
    current: list[str] = []
    current_len = 0
    page: int | None = None
    for p in paragraphs:
        p = p.strip()
        if not p:
            continue
        if current_len + len(p) + 2 <= max_chars:
            current.append(p)
            current_len += len(p) + 2
        else:
            if current:
                chunks.append(("\n\n".join(current), page))
            # Start new chunk; optionally carry over overlap
            if overlap and current:
                overlap_text = "\n\n".join(current)[-overlap:].lstrip()
                current = [overlap_text] if overlap_text else []
                current_len = len(overlap_text)
            else:
                current = []
                current_len = 0
            # Fit as much of p as possible
            while len(p) > max_chars:
                chunks.append((p[:max_chars], page))
                p = p[max_chars - overlap :].lstrip()
            if p:
                current.append(p)
                current_len = len(p) + 2
    if current:
        chunks.append(("\n\n".join(current), page))
    return chunks


def parse_pdf(path: Path) -> list[tuple[str, int | None]]:
    import fitz  # PyMuPDF
    chunks: list[tuple[str, int | None]] = []
    doc = fitz.open(path)
    try:
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if not text.strip():
                continue
            page_chunks = chunk_text(text, max_chars=1200, overlap=100)
            for content, _ in page_chunks:
                chunks.append((content, page_num + 1))
    finally:
        doc.close()
    return chunks


def parse_txt(path: Path) -> list[tuple[str, int | None]]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return chunk_text(text, max_chars=1200, overlap=100)


def parse_pptx(path: Path) -> list[tuple[str, int | None]]:
    from pptx import Presentation
    from pptx.util import Inches, Pt  # noqa: F401

    prs = Presentation(path)
    chunks: list[tuple[str, int | None]] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text.strip())
        text = "\n\n".join(p for p in parts if p)
        if not text:
            continue
        slide_chunks = chunk_text(text, max_chars=1200, overlap=100)
        for content, _ in slide_chunks:
            chunks.append((content, slide_num))
    return chunks


def parse_file(path: Path, file_type: str) -> list[tuple[str, int | None]]:
    if file_type == "pdf":
        return parse_pdf(path)
    if file_type == "txt":
        return parse_txt(path)
    if file_type == "pptx":
        return parse_pptx(path)
    raise ValueError(f"Unsupported file type: {file_type}")
